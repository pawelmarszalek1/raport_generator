from pptx import Presentation
from pptx.util import Inches, Pt
import re
import io

def _add_bullet_slide(prs, layout, title_text, content_text):
    """Tworzy nowy slajd i wypełnia go punktami (bullet points)."""
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text
    
    content_box = slide.shapes.placeholders[1]
    tf = content_box.text_frame
    tf.clear()  # Czyści domyślny tekst

    # Dzieli tekst od Gemini na linie (punkty)
    lines = content_text.split('\n')
    
    if not lines or (len(lines) == 1 and not lines[0]):
        p = tf.add_paragraph()
        p.text = "Brak danych."
        return

    # Dodaje pierwszą linię jako pierwszy paragraf
    p = tf.paragraphs[0]
    p.text = lines[0].strip()
    
    # Dodaje resztę linii jako nowe paragrafy (punkty)
    for line in lines[1:]:
        if line.strip(): # Pomiń puste linie
            p = tf.add_paragraph()
            p.text = line.strip()

def parse_gemini_response(raw_text):
    """Dzieli surową odpowiedź Gemini na słownik sekcji."""
    sections = {}
    # Używamy wyrażeń regularnych do znalezienia naszych znaczników
    pattern = re.compile(r"\[SEKCJA:([\w_]+)\]\n(.*?)(?=\n\[SEKCJA:|\Z)", re.DOTALL)
    
    matches = pattern.findall(raw_text)
    
    if not matches:
        # Jeśli Gemini nie zwróciło tagów, zwróć wszystko jako jedną sekcję
        return {'BRAK_TAGOW': raw_text}

    for name, content in matches:
        sections[name] = content.strip()
        
    return sections

# --- ZAKTUALIZOWANA DEFINICJA FUNKCJI ---
def create_pptx_report(gemini_text, template_path, report_title, report_subtitle):
    """
    Tworzy raport PPTX na podstawie treści od Gemini i szablonu.
    Uzupełnia również slajd tytułowy.
    """
    
    insights = parse_gemini_response(gemini_text)
    
    # Otwórz szablon prezentacji
    prs = Presentation(template_path)

    # --- NOWA SEKCJA: Edycja slajdu tytułowego ---
    # Zakładamy, że pierwszy slajd (index 0) w szablonie to slajd tytułowy
    try:
        title_slide = prs.slides[0]
        
        # Ustawiamy główny tytuł (zazwyczaj 'shapes.title')
        if title_slide.shapes.title:
            title_slide.shapes.title.text = report_title
        
        # Ustawiamy podtytuł (zazwyczaj placeholder o idx 1)
        # Sprawdzamy, czy placeholder istnieje i jest to placeholder
        if len(title_slide.placeholders) > 1:
             subtitle_placeholder = title_slide.placeholders[1]
             if subtitle_placeholder:
                subtitle_placeholder.text = report_subtitle
    except Exception as e:
        # Drukujemy ostrzeżenie w konsoli, ale nie przerywamy generowania raportu
        print(f"Ostrzeżenie: Nie można było ustawić tytułu/podtytułu na slajdzie 0. Upewnij się, że szablon ma slajd tytułowy. Błąd: {e}")
    # --- KONIEC NOWEJ SEKCJI ---


    # === Slajd 2: Podsumowanie dla Zarządu ===
    # Zakładamy, że layout [1] w szablonie to 'Tytuł i Treść'
    # Jeśli Twój szablon ma inaczej, musisz zmienić ten indeks
    # Znajdź właściwy layout slajdu (zazwyczaj indeks 1)
    try:
        content_layout = prs.slide_layouts[1] 
    except IndexError:
        content_layout = prs.slide_layouts[1] # Na wypadek błędu szablonu

    # Używamy naszej nowej funkcji pomocniczej
    _add_bullet_slide(
        prs, 
        content_layout, 
        "Podsumowanie dla Zarządu", 
        insights.get('PODSUMOWANIE_ZARZADU', 'Brak danych.')
    )
    
    _add_bullet_slide(
        prs, 
        content_layout, 
        "Analiza Sprzedaży", 
        insights.get('ANALIZA_SPRZEDAZY', 'Brak danych.')
    )
    
    _add_bullet_slide(
        prs, 
        content_layout, 
        "Analiza Marketingu i KPI", 
        insights.get('ANALIZA_MARKETINGU', 'Brak danych.')
    )
    
    _add_bullet_slide(
        prs, 
        content_layout, 
        "Wnioski i Rekomendacje", 
        insights.get('WNIOSKI_I_REKOMENDACJE', 'Brak danych.')
    )

    # TODO: W tym miejscu można dodać slajd z wykresem wygenerowanym przez Matplotlib

    # Zapisz prezentację do bufora w pamięci, aby Streamlit mógł ją udostępnić
    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    
    return file_stream